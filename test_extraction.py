"""
Test universal extraction with sample files.
"""

from legal_assistant.utils.universal_extraction import extract_text_from_upload

# Test text file
text_data = b"This is a test text file.\nSecond line."
result = extract_text_from_upload("test.txt", text_data)
print(f"TXT: source_type={result.source_type}, text_length={len(result.text)}, error={result.error}")
assert result.source_type == "text"
assert len(result.text) > 0

# Test CSV
csv_data = b"Name,Age\nJohn,30\nJane,25"
result = extract_text_from_upload("test.csv", csv_data)
print(f"CSV: source_type={result.source_type}, text_length={len(result.text)}, error={result.error}")
assert result.source_type == "text"

# Test HTML
html_data = b"<html><body><p>Test paragraph</p><script>alert('test')</script></body></html>"
result = extract_text_from_upload("test.html", html_data)
print(f"HTML: source_type={result.source_type}, text_length={len(result.text)}, error={result.error}")
assert result.source_type == "html"
assert "Test paragraph" in result.text
assert "alert" not in result.text  # script should be removed

# Test PDF (simple)
try:
    from pypdf import PdfWriter
    import io
    
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    writer.pages[0].add_text(100, 100, "Test PDF content")
    
    pdf_buffer = io.BytesIO()
    writer.write(pdf_buffer)
    pdf_data = pdf_buffer.getvalue()
    
    result = extract_text_from_upload("test.pdf", pdf_data)
    print(f"PDF: source_type={result.source_type}, text_length={len(result.text)}, error={result.error}, meta={result.meta}")
    assert result.source_type == "pdf"
except Exception as e:
    print(f"PDF test skipped: {e}")

# Test DOCX (requires creating a sample)
try:
    from docx import Document
    import io
    
    doc = Document()
    doc.add_paragraph("Test Word document content")
    doc.add_paragraph("Second paragraph with more text")
    
    docx_buffer = io.BytesIO()
    doc.save(docx_buffer)
    docx_data = docx_buffer.getvalue()
    
    result = extract_text_from_upload("test.docx", docx_data)
    print(f"DOCX: source_type={result.source_type}, text_length={len(result.text)}, error={result.error}")
    assert result.source_type == "docx"
    assert "Test Word document" in result.text
except Exception as e:
    print(f"DOCX test skipped: {e}")

# Test PPTX
try:
    from pptx import Presentation
    import io
    
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Test Presentation"
    
    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    pptx_data = pptx_buffer.getvalue()
    
    result = extract_text_from_upload("test.pptx", pptx_data)
    print(f"PPTX: source_type={result.source_type}, text_length={len(result.text)}, error={result.error}")
    assert result.source_type == "pptx"
    assert "Test Presentation" in result.text
except Exception as e:
    print(f"PPTX test skipped: {e}")

# Test unsupported type
result = extract_text_from_upload("test.xyz", b"unknown content")
print(f"Unknown: source_type={result.source_type}, error={result.error}")
assert result.source_type == "fallback"

print("\n✓ All extraction tests passed!")

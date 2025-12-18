# legal_assistant/utils/universal_extraction.py
"""
Universal text extraction module.
Supports: PDF, DOCX, PPTX, HTML, TXT, MD, CSV, LOG, EML, and audio files.
"""
from __future__ import annotations

import os
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

# Existing EML extraction
from legal_assistant.utils.eml_extraction import extract_eml_from_bytes


@dataclass
class ExtractedText:
    """Result of text extraction from an uploaded file."""
    text: str
    source_type: str
    meta: dict = field(default_factory=dict)
    error: Optional[str] = None


# ---------------------------------------------------------------------------
# Text-based files (TXT, MD, CSV, LOG)
# ---------------------------------------------------------------------------

TEXT_EXTENSIONS = {".txt", ".md", ".csv", ".log"}


def _extract_text_file(data: bytes) -> ExtractedText:
    """Decode UTF-8 text files."""
    try:
        text = data.decode("utf-8", errors="ignore")
        return ExtractedText(text=text, source_type="text", meta={})
    except Exception as e:
        return ExtractedText(text="", source_type="text", error=f"Text decode error: {e}")


# ---------------------------------------------------------------------------
# EML files
# ---------------------------------------------------------------------------

def _extract_eml(data: bytes) -> ExtractedText:
    """Extract email content using existing EML parser."""
    try:
        eml_data = extract_eml_from_bytes(data)
        if not eml_data:
            return ExtractedText(text="", source_type="eml", error="No usable content in email")
        
        header = (
            f"From: {eml_data['from']}\n"
            f"To: {eml_data['to']}\n"
            f"Subject: {eml_data['subject']}\n"
            f"Date: {eml_data['date']}\n\n"
        )
        text = header + eml_data["body"]
        return ExtractedText(
            text=text,
            source_type="eml",
            meta={
                "from": eml_data["from"],
                "to": eml_data["to"],
                "subject": eml_data["subject"],
                "date": eml_data["date"],
            }
        )
    except Exception as e:
        return ExtractedText(text="", source_type="eml", error=f"EML parse error: {e}")


# ---------------------------------------------------------------------------
# Document files (PDF, DOCX, PPTX, HTML)
# ---------------------------------------------------------------------------

DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".pptx", ".html", ".htm"}


def _extract_pdf(data: bytes) -> ExtractedText:
    """Extract text from PDF using pypdf."""
    try:
        from pypdf import PdfReader
        import io
        
        reader = PdfReader(io.BytesIO(data))
        texts = []
        for page in reader.pages:
            texts.append(page.extract_text() or "")
        
        text = "\n".join(texts).strip()
        return ExtractedText(
            text=text,
            source_type="pdf",
            meta={"pages": len(reader.pages)}
        )
    except Exception as e:
        return ExtractedText(
            text="",
            source_type="pdf",
            error=f"PDF extraction error: {e}"
        )


def _extract_docx(data: bytes) -> ExtractedText:
    """Extract text from DOCX using python-docx."""
    try:
        from docx import Document
        import io
        
        doc = Document(io.BytesIO(data))
        texts = [para.text for para in doc.paragraphs if para.text.strip()]
        text = "\n\n".join(texts)
        
        return ExtractedText(
            text=text,
            source_type="docx",
            meta={"paragraphs": len(doc.paragraphs)}
        )
    except ImportError:
        return ExtractedText(
            text="",
            source_type="docx",
            error="python-docx not installed. Run: pip install python-docx"
        )
    except Exception as e:
        return ExtractedText(
            text="",
            source_type="docx",
            error=f"DOCX extraction error: {e}"
        )


def _extract_pptx(data: bytes) -> ExtractedText:
    """Extract text from PPTX using python-pptx."""
    try:
        from pptx import Presentation
        import io
        
        prs = Presentation(io.BytesIO(data))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        
        text = "\n\n".join(texts)
        return ExtractedText(
            text=text,
            source_type="pptx",
            meta={"slides": len(prs.slides)}
        )
    except ImportError:
        return ExtractedText(
            text="",
            source_type="pptx",
            error="python-pptx not installed. Run: pip install python-pptx"
        )
    except Exception as e:
        return ExtractedText(
            text="",
            source_type="pptx",
            error=f"PPTX extraction error: {e}"
        )


def _extract_html(data: bytes) -> ExtractedText:
    """Extract text from HTML using BeautifulSoup."""
    try:
        from bs4 import BeautifulSoup
        
        text = data.decode("utf-8", errors="ignore")
        soup = BeautifulSoup(text, "html.parser")
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        
        text = soup.get_text(separator="\n", strip=True)
        return ExtractedText(
            text=text,
            source_type="html",
            meta={}
        )
    except ImportError:
        return ExtractedText(
            text="",
            source_type="html",
            error="beautifulsoup4 not installed. Run: pip install beautifulsoup4"
        )
    except Exception as e:
        return ExtractedText(
            text="",
            source_type="html",
            error=f"HTML extraction error: {e}"
        )


def _extract_document(filename: str, data: bytes) -> ExtractedText:
    """Route document extraction based on file type."""
    ext = Path(filename).suffix.lower()
    
    if ext == ".pdf":
        return _extract_pdf(data)
    elif ext == ".docx":
        return _extract_docx(data)
    elif ext == ".pptx":
        return _extract_pptx(data)
    elif ext in {".html", ".htm"}:
        return _extract_html(data)
    else:
        return ExtractedText(
            text="",
            source_type="unknown",
            error=f"Unsupported document type: {ext}"
        )


# ---------------------------------------------------------------------------
# Audio files via Whisper transcription
# ---------------------------------------------------------------------------

AUDIO_EXTENSIONS = {".mp3", ".wav", ".m4a", ".aac", ".flac", ".ogg", ".wma", ".webm"}


def _transcribe_audio_azure(filename: str, data: bytes) -> ExtractedText:
    """
    Transcribe audio using Azure OpenAI Whisper API.
    Falls back to local faster-whisper if Azure is unavailable.
    """
    ext = Path(filename).suffix.lower()
    
    # Try Azure OpenAI first
    try:
        from openai import AzureOpenAI
        from legal_assistant.config import get_settings
        
        settings = get_settings()
        
        base_url = os.getenv("OPENAI_BASE_URL")
        api_key = os.getenv("OPENAI_API_KEY")
        api_version = os.getenv("OPENAI_API_VERSION")
        
        if base_url and api_key and api_version:
            client = AzureOpenAI(
                azure_endpoint=base_url,
                api_key=api_key,
                api_version=api_version,
            )
            
            # Write to temp file for API
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                tmp.write(data)
                tmp_path = tmp.name
            
            try:
                with open(tmp_path, "rb") as audio_file:
                    # Use whisper-1 model for transcription
                    whisper_model = os.getenv("OPENAI_WHISPER_MODEL", "whisper-1")
                    transcription = client.audio.transcriptions.create(
                        model=whisper_model,
                        file=audio_file,
                        response_format="text"
                    )
                
                text = transcription if isinstance(transcription, str) else str(transcription)
                
                print(f"[EXTRACT] Audio transcribed via Azure Whisper: {filename}")
                return ExtractedText(
                    text=text,
                    source_type="audio",
                    meta={"transcription_method": "azure_whisper", "format": ext.lstrip(".")}
                )
            finally:
                try:
                    os.unlink(tmp_path)
                except OSError:
                    pass
    except Exception as azure_err:
        print(f"[WARN] Azure Whisper unavailable: {azure_err}, trying local faster-whisper")

    # Fallback to local faster-whisper
    return _transcribe_audio_local(filename, data)


def _transcribe_audio_local(filename: str, data: bytes) -> ExtractedText:
    """
    Transcribe audio using local OpenAI Whisper.
    Requires ffmpeg to be installed and in PATH.
    """
    ext = Path(filename).suffix.lower()
    
    try:
        import whisper
    except ImportError:
        return ExtractedText(
            text="",
            source_type="audio",
            error="Audio transcription requires openai-whisper. Run: pip install openai-whisper"
        )
    
    # Check for ffmpeg
    import shutil
    if not shutil.which("ffmpeg"):
        return ExtractedText(
            text="",
            source_type="audio",
            error="Audio transcription requires ffmpeg. Install from https://ffmpeg.org/download.html or run: winget install ffmpeg"
        )

    # Write to temp file (must close file before Whisper can read it on Windows)
    tmp_file = tempfile.NamedTemporaryFile(delete=False, suffix=ext)
    tmp_path = tmp_file.name
    
    try:
        tmp_file.write(data)
        tmp_file.flush()
        tmp_file.close()  # Must close before Whisper can read on Windows
        
        # Verify file exists and has content
        if not os.path.exists(tmp_path):
            raise FileNotFoundError(f"Temp file not created: {tmp_path}")
        if os.path.getsize(tmp_path) == 0:
            raise ValueError("Temp audio file is empty")
        
        # Use base model for balance of speed/accuracy; can be configured
        model_size = os.getenv("WHISPER_MODEL_SIZE", "base")
        model = whisper.load_model(model_size)
        
        result = model.transcribe(tmp_path)
        
        text = result["text"].strip()
        
        print(f"[EXTRACT] Audio transcribed via local OpenAI Whisper: {filename}")
        return ExtractedText(
            text=text,
            source_type="audio",
            meta={
                "transcription_method": "openai_whisper",
                "format": ext.lstrip("."),
                "language": result.get("language", "unknown"),
            }
        )
    except Exception as e:
        return ExtractedText(
            text="",
            source_type="audio",
            error=f"Audio transcription error: {e}"
        )
    finally:
        try:
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)
        except OSError:
            pass


# ---------------------------------------------------------------------------
# Main extraction router
# ---------------------------------------------------------------------------

def extract_text_from_upload(filename: str, data: bytes) -> ExtractedText:
    """
    Universal text extraction from uploaded file bytes.
    
    Routes to appropriate extractor based on file extension:
    - .txt, .md, .csv, .log → UTF-8 decode
    - .eml → Email parser
    - .pdf → pypdf extraction
    - .docx → python-docx extraction
    - .pptx → python-pptx extraction
    - .html, .htm → BeautifulSoup extraction
    - .mp3, .wav, .m4a, .aac, .flac, .ogg, .wma, .webm → Whisper transcription
    - Unknown → Fallback UTF-8 decode
    """
    if not data:
        return ExtractedText(text="", source_type="unknown", error="Empty file")

    ext = Path(filename).suffix.lower()
    
    # Text files
    if ext in TEXT_EXTENSIONS:
        print(f"[EXTRACT] Text extraction: {filename}")
        return _extract_text_file(data)
    
    # Email files
    if ext == ".eml":
        print(f"[EXTRACT] EML extraction: {filename}")
        return _extract_eml(data)
    
    # Document files (PDF, DOCX, PPTX, HTML)
    if ext in DOCUMENT_EXTENSIONS:
        print(f"[EXTRACT] Document extraction: {filename}")
        return _extract_document(filename, data)
    
    # Audio files
    if ext in AUDIO_EXTENSIONS:
        print(f"[EXTRACT] Audio transcription: {filename}")
        return _transcribe_audio_azure(filename, data)
    
    # Unknown extension - try text fallback
    print(f"[EXTRACT] Fallback text decode: {filename}")
    result = _extract_text_file(data)
    result.source_type = "fallback"
    return result


def get_supported_extensions() -> dict:
    """Return dict of supported file types and their extensions."""
    return {
        "text": list(TEXT_EXTENSIONS),
        "email": [".eml"],
        "document": list(DOCUMENT_EXTENSIONS),
        "audio": list(AUDIO_EXTENSIONS),
    }
